import difflib
import io
import logging
import os
import re
import shutil
import shlex
import tempfile
import ast
import operator
import subprocess
import zipfile
from html import escape
from copy import deepcopy
from pathlib import Path
from typing import List

from flask import Flask, abort, jsonify, redirect, request, send_file
from flask_cors import CORS
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
from pypdf import PdfReader, PdfWriter
from docx import Document
from docx.shared import Inches
from openpyxl import load_workbook, Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import ParagraphStyle
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
import reportlab
from werkzeug.utils import secure_filename
from werkzeug.exceptions import HTTPException
from werkzeug.middleware.proxy_fix import ProxyFix

logging.basicConfig(level=logging.INFO)

app = Flask(__name__, static_folder=".", static_url_path="")
app.config["MAX_CONTENT_LENGTH"] = 16 * 1024 * 1024
app.wsgi_app = ProxyFix(app.wsgi_app, x_for=1, x_proto=1, x_host=1)
allowed_origins = [
    origin.strip()
    for origin in os.environ.get(
        "ALLOWED_ORIGINS",
        "https://belgelab.com.tr,https://www.belgelab.com.tr,http://127.0.0.1:3000,http://localhost:3000,capacitor://localhost",
    ).split(",")
    if origin.strip()
]
CORS(app, resources={r"/api/*": {"origins": allowed_origins}}, methods=["POST", "OPTIONS"])
limiter = Limiter(
    key_func=get_remote_address,
    app=app,
    default_limits=[],
    storage_uri=os.environ.get("RATELIMIT_STORAGE_URI", "memory://"),
)

SUPPORTED_OPERATIONS = {
    "pdf-to-docx": {".pdf"},
    "pdf-to-xlsx": {".pdf"},
    "pdf-to-pptx": {".pdf"},
    "docx-to-pdf": {".docx"},
    "xlsx-to-pdf": {".xlsx"},
    "pptx-to-pdf": {".pptx"},
    "dwg-to-pdf": {".dwg"},
}
APP_VERSION = "2.3.0"
INVALID_XML_CHARS = re.compile(r"[\x00-\x08\x0B\x0C\x0E-\x1F\x7F-\x9F]")
MAX_PDF_PAGES = 250
MAX_ARCHIVE_FILES = 10_000
MAX_ARCHIVE_EXPANDED_BYTES = 64 * 1024 * 1024
OFFICE_MARKERS = {
    ".docx": "word/document.xml",
    ".xlsx": "xl/workbook.xml",
    ".pptx": "ppt/presentation.xml",
}


class ConversionServiceUnavailable(RuntimeError):
    pass


def clean_document_text(value: str) -> str:
    return INVALID_XML_CHARS.sub("", value or "")


def validate_pdf_bytes(data: bytes) -> None:
    if not data or b"%PDF-" not in data[:1024]:
        raise ValueError("Dosyanın gerçek bir PDF olduğu doğrulanamadı.")


def validate_office_archive(path: Path, extension: str) -> None:
    marker = OFFICE_MARKERS.get(extension)
    if not marker or not zipfile.is_zipfile(path):
        raise ValueError("Dosyanın gerçek bir Office belgesi olduğu doğrulanamadı.")
    try:
        with zipfile.ZipFile(path) as archive:
            entries = archive.infolist()
            if len(entries) > MAX_ARCHIVE_FILES:
                raise ValueError("Belge çok fazla iç dosya içeriyor.")
            expanded_size = sum(entry.file_size for entry in entries)
            if expanded_size > MAX_ARCHIVE_EXPANDED_BYTES:
                raise ValueError("Belgenin açılmış boyutu güvenlik sınırını aşıyor.")
            names = {entry.filename for entry in entries}
            if "[Content_Types].xml" not in names or marker not in names:
                raise ValueError("Office belge yapısı geçersiz.")
    except zipfile.BadZipFile as exc:
        raise ValueError("Office arşivi okunamadı.") from exc


def validate_dwg_file(path: Path) -> None:
    signature = path.read_bytes()[:6]
    if len(signature) != 6 or not re.fullmatch(rb"AC10\d{2}", signature):
        raise ValueError("Dosyanın gerçek bir DWG çizimi olduğu doğrulanamadı.")


def enforce_pdf_page_limit(reader: PdfReader) -> None:
    if len(reader.pages) > MAX_PDF_PAGES:
        raise ValueError(f"PDF en fazla {MAX_PDF_PAGES} sayfa içerebilir.")


@app.before_request
def enforce_canonical_origin():
    protected_paths = {
        "/server.py",
        "/requirements.txt",
        "/Dockerfile",
        "/docker-compose.yml",
        "/.dockerignore",
    }
    if request.path in protected_paths or request.path.startswith("/tessdata/"):
        abort(404)
    hostname = request.host.split(":", 1)[0].lower()
    if hostname in {"belgelab.com.tr", "www.belgelab.com.tr"}:
        if not request.is_secure or hostname != "belgelab.com.tr":
            canonical_url = f"https://belgelab.com.tr{request.full_path}"
            if canonical_url.endswith("?"):
                canonical_url = canonical_url[:-1]
            return redirect(canonical_url, code=308)
    return None


@app.after_request
def disable_stale_app_cache(response):
    if request.path.startswith("/api/"):
        response.headers["Cache-Control"] = "no-store, max-age=0"
        response.headers["Pragma"] = "no-cache"
        response.headers["Expires"] = "0"
        response.headers["Clear-Site-Data"] = '"cache"'
    elif request.path in {"/", "/index.html", "/sw.js", "/app.js", "/tools.js", "/style.css"}:
        response.headers["Cache-Control"] = "no-cache, no-store, must-revalidate"
        response.headers["Pragma"] = "no-cache"
        response.headers["Expires"] = "0"
    response.headers["X-Content-Type-Options"] = "nosniff"
    response.headers["Referrer-Policy"] = "strict-origin-when-cross-origin"
    response.headers["X-Frame-Options"] = "DENY"
    response.headers["Permissions-Policy"] = "camera=(), microphone=(), geolocation=(), payment=(), usb=()"
    response.headers["Cross-Origin-Opener-Policy"] = "same-origin-allow-popups"
    response.headers["Strict-Transport-Security"] = "max-age=31536000; includeSubDomains"
    response.headers["Content-Security-Policy"] = (
        "default-src 'self'; "
        "base-uri 'self'; object-src 'none'; frame-ancestors 'none'; form-action 'self'; "
        "script-src 'self' 'sha256-kwaLn3eP/r2tN51WK5mNfToQKOUTowmq0lYDbpUQ9Qs=' https://cdnjs.cloudflare.com "
        "https://pagead2.googlesyndication.com https://*.googlesyndication.com https://*.google.com https://*.gstatic.com; "
        "style-src 'self' 'unsafe-inline'; "
        "img-src 'self' data: blob: https://*.google.com https://*.gstatic.com https://*.googlesyndication.com; "
        "connect-src 'self' https://cdnjs.cloudflare.com https://*.google.com https://*.googlesyndication.com; "
        "worker-src 'self' blob: https://cdnjs.cloudflare.com; "
        "frame-src https://*.google.com https://*.googlesyndication.com; "
        "upgrade-insecure-requests"
    )
    return response


@app.errorhandler(413)
def file_too_large(_error):
    return jsonify({"error": "Dosya çok büyük. En fazla 16 MB yükleyebilirsiniz."}), 413


@app.errorhandler(429)
def rate_limited(_error):
    return jsonify({"error": "Çok fazla istek gönderdiniz. Lütfen kısa bir süre sonra yeniden deneyin."}), 429


@app.errorhandler(Exception)
def unhandled_error(error):
    if isinstance(error, HTTPException):
        return error
    logging.exception("Unhandled server error")
    return jsonify({"error": "Sunucu işlemi tamamlayamadı. Lütfen dosyayı kontrol edip yeniden deneyin."}), 500


def extract_text_from_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    pages: List[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        pages.append(text.strip())
    return "\n\n".join(page for page in pages if page)


def extract_text_from_docx(path: Path) -> str:
    doc = Document(str(path))
    return "\n".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text)


def extract_text_from_xlsx(path: Path) -> str:
    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    sheet = workbook.active
    rows = []
    for row in sheet.iter_rows(values_only=True):
        cleaned = ["" if value is None else str(value) for value in row]
        rows.append(" | ".join(cleaned))
    workbook.close()
    return "\n".join(rows)


def extract_text_from_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    slides_text: List[str] = []
    for slide in presentation.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        slides_text.append("\n".join(texts) if texts else "(Boş slayt)")
    return "\n\n".join(slides_text)


def create_docx_from_text(text: str, output_path: Path) -> None:
    doc = Document()
    for paragraph in clean_document_text(text).splitlines():
        doc.add_paragraph(paragraph or "")
    doc.save(str(output_path))


def create_docx_from_scanned_pdf(input_path: Path, output_path: Path) -> None:
    import pypdfium2 as pdfium

    pdf = pdfium.PdfDocument(str(input_path))
    try:
        doc = Document()
        section = doc.sections[0]
        section.top_margin = Inches(0.35)
        section.bottom_margin = Inches(0.35)
        section.left_margin = Inches(0.35)
        section.right_margin = Inches(0.35)
        usable_width = section.page_width - section.left_margin - section.right_margin
        page_count = len(pdf)

        for index in range(page_count):
            page = pdf[index]
            try:
                bitmap = page.render(scale=2.0)
                try:
                    image = bitmap.to_pil().convert("RGB")
                finally:
                    bitmap.close()
            finally:
                page.close()
            image_stream = io.BytesIO()
            image.save(image_stream, format="JPEG", quality=92, optimize=True)
            image.close()
            image_stream.seek(0)
            paragraph = doc.add_paragraph()
            paragraph.alignment = 1
            run = paragraph.add_run()
            run.add_picture(image_stream, width=usable_width)
            image_stream.close()
            if index < page_count - 1:
                doc.add_page_break()

        doc.save(str(output_path))
    finally:
        pdf.close()


def find_tesseract_executable() -> str:
    candidates = [
        shutil.which("tesseract"),
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return str(candidate)
    raise RuntimeError("OCR motoru bulunamadı. Tesseract OCR kurulmalıdır.")


def create_editable_docx_with_ocr(input_path: Path, output_path: Path) -> bool:
    import pypdfium2 as pdfium
    import pytesseract

    pytesseract.pytesseract.tesseract_cmd = find_tesseract_executable()
    bundled_tessdata = Path(__file__).resolve().parent / "tessdata" / "tur.traineddata"
    if not bundled_tessdata.exists():
        raise RuntimeError("Türkçe OCR dil modeli bulunamadı.")
    tessdata_dir = Path(tempfile.gettempdir()) / "belgelab-tessdata"
    tessdata_dir.mkdir(parents=True, exist_ok=True)
    runtime_language_file = tessdata_dir / "tur.traineddata"
    if not runtime_language_file.exists() or runtime_language_file.stat().st_size != bundled_tessdata.stat().st_size:
        shutil.copy2(bundled_tessdata, runtime_language_file)

    pdf = pdfium.PdfDocument(str(input_path))
    doc = Document()
    recognized_any_text = False
    try:
        for index in range(len(pdf)):
            page = pdf[index]
            try:
                bitmap = page.render(scale=3.0)
                try:
                    image = bitmap.to_pil().convert("L")
                finally:
                    bitmap.close()
            finally:
                page.close()

            text = pytesseract.image_to_string(
                image,
                lang="tur",
                config=f"--tessdata-dir {tessdata_dir} --psm 6",
            )
            image.close()
            lines = [clean_document_text(line).strip() for line in text.splitlines() if line.strip()]
            if lines:
                recognized_any_text = True
                for line in lines:
                    paragraph = doc.add_paragraph(line)
                    paragraph.paragraph_format.space_after = Inches(0.04)
            if index < len(pdf) - 1:
                doc.add_page_break()

        if recognized_any_text:
            doc.save(str(output_path))
        return recognized_any_text
    finally:
        pdf.close()


def create_hybrid_docx(input_path: Path, output_path: Path) -> None:
    create_docx_from_scanned_pdf(input_path, output_path)
    with tempfile.TemporaryDirectory() as temp_dir:
        ocr_path = Path(temp_dir) / "ocr.docx"
        if not create_editable_docx_with_ocr(input_path, ocr_path):
            return
        result = Document(str(output_path))
        result.add_page_break()
        result.add_heading("Düzenlenebilir OCR Metni", level=1)
        ocr_document = Document(str(ocr_path))
        for element in ocr_document.element.body:
            if not element.tag.endswith("sectPr"):
                result.element.body.append(deepcopy(element))
        result.save(str(output_path))


def create_xlsx_from_text(text: str, output_path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Dönüştürülmüş"
    for line in clean_document_text(text).splitlines() or [""]:
        sheet.append([line[:32767]])
    workbook.save(str(output_path))


def parse_turkish_number(value: str):
    cleaned = (value or "").strip()
    if not cleaned:
        return None
    if re.fullmatch(r"-?\d{1,3}(?:\.\d{3})*,\d+", cleaned):
        return float(cleaned.replace(".", "").replace(",", "."))
    if re.fullmatch(r"-?\d+,\d+", cleaned):
        return float(cleaned.replace(",", "."))
    if re.fullmatch(r"-?\d+", cleaned):
        return int(cleaned)
    return cleaned


def create_xlsx_from_pdf_tables(input_path: Path, output_path: Path) -> bool:
    try:
        import pdfplumber
    except ImportError:
        return False

    extracted_tables = []
    with pdfplumber.open(str(input_path)) as pdf:
        for page in pdf.pages:
            tables = page.extract_tables()
            extracted_tables.extend(table for table in tables if table)
    if not extracted_tables:
        return False

    first_table = extracted_tables[0]
    header_index = next(
        (index for index, row in enumerate(first_table)
         if row and any("SIRA" in clean_document_text(cell or "").upper() for cell in row)),
        None,
    )
    if header_index is None:
        return False

    raw_headers = first_table[header_index]
    group_row = first_table[header_index - 1] if header_index > 0 else [None] * len(raw_headers)
    headers = [
        clean_document_text(raw_header or f"Sütun {index + 1}").replace("\n", " ").strip()
        for index, raw_header in enumerate(raw_headers)
    ]
    group_headers = [
        clean_document_text(value or "").replace("\n", " ").strip()
        for value in group_row
    ]

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "PDF Tablosu"
    sheet.append(group_headers)
    sheet.append(headers)
    numeric_columns = set(range(9, len(headers)))

    for table in extracted_tables:
        page_header_index = next(
            (index for index, row in enumerate(table)
             if row and any("SIRA" in clean_document_text(cell or "").upper() for cell in row)),
            None,
        )
        data_start = page_header_index + 1 if page_header_index is not None else 0
        for row in table[data_start:]:
            if not row or not any(cell not in (None, "") for cell in row):
                continue
            first_cell = clean_document_text(row[0] or "").strip()
            if not re.fullmatch(r"\d+", first_cell):
                continue
            normalized = []
            for index in range(len(headers)):
                value = clean_document_text(row[index] or "").replace("\n", " ").strip() if index < len(row) else ""
                normalized.append(parse_turkish_number(value) if index in numeric_columns else value)
            sheet.append(normalized)

    if sheet.max_row == 2:
        return False

    group_starts = [index for index, value in enumerate(group_headers) if value]
    for position, start in enumerate(group_starts):
        end = group_starts[position + 1] - 1 if position + 1 < len(group_starts) else len(headers) - 1
        if end > start:
            sheet.merge_cells(start_row=1, start_column=start + 1, end_row=1, end_column=end + 1)

    group_fill = PatternFill("solid", fgColor="252039")
    for cell in sheet[1]:
        cell.fill = group_fill
        cell.font = Font(color="FFFFFF", bold=True)
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    header_fill = PatternFill("solid", fgColor="6C4EE3")
    for cell in sheet[2]:
        cell.fill = header_fill
        cell.font = Font(color="FFFFFF", bold=True)
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    sheet.row_dimensions[1].height = 42
    sheet.row_dimensions[2].height = 52
    sheet.freeze_panes = "A3"
    sheet.auto_filter.ref = f"A2:{get_column_letter(sheet.max_column)}{sheet.max_row}"

    for column_index in range(1, sheet.max_column + 1):
        letter = get_column_letter(column_index)
        values = [str(sheet.cell(row=row, column=column_index).value or "") for row in range(1, min(sheet.max_row, 80) + 1)]
        sheet.column_dimensions[letter].width = min(28, max(10, max(map(len, values), default=10) + 2))
        if column_index - 1 in numeric_columns:
            for cell in sheet[letter][2:]:
                if isinstance(cell.value, float):
                    cell.number_format = '#,##0.00'
                elif isinstance(cell.value, int):
                    cell.number_format = '#,##0'

    workbook.save(str(output_path))
    return True


def create_pptx_from_text(text: str, output_path: Path) -> None:
    presentation = Presentation()
    for chunk in clean_document_text(text).splitlines() or [""]:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "İçerik"
        slide.placeholders[1].text = chunk or "(Boş içerik)"
    presentation.save(str(output_path))


def create_pdf_from_text(text: str, output_path: Path) -> None:
    pdf = canvas.Canvas(str(output_path), pagesize=letter)
    width, height = letter
    font_name = register_unicode_pdf_font()
    y = height - 40
    lines = clean_document_text(text).splitlines() or [""]
    for line in lines:
        wrapped = wrap_pdf_line(line, font_name, 10, width - 80)
        for wrapped_line in wrapped:
            if y < 40:
                pdf.showPage()
                pdf.setFont(font_name, 10)
                y = height - 40
            pdf.setFont(font_name, 10)
            pdf.drawString(40, y, wrapped_line)
            y -= 14
    pdf.save()


def register_unicode_pdf_font() -> str:
    font_name = "BelgeAtolyeUnicode"
    if font_name not in pdfmetrics.getRegisteredFontNames():
        fonts_dir = Path(reportlab.__file__).resolve().parent / "fonts"
        font_path = fonts_dir / "Vera.ttf"
        bold_path = fonts_dir / "VeraBd.ttf"
        if not font_path.exists():
            raise RuntimeError("Unicode PDF yazı tipi bulunamadı.")
        pdfmetrics.registerFont(TTFont(font_name, str(font_path)))
        if bold_path.exists():
            pdfmetrics.registerFont(TTFont(f"{font_name}-Bold", str(bold_path)))
            pdfmetrics.registerFontFamily(
                font_name,
                normal=font_name,
                bold=f"{font_name}-Bold",
                italic=font_name,
                boldItalic=f"{font_name}-Bold",
            )
    return font_name


def wrap_pdf_line(text: str, font_name: str, font_size: int, max_width: float) -> List[str]:
    if not text:
        return [""]
    words = text.split()
    if not words:
        return [""]
    lines = []
    current = words[0]
    for word in words[1:]:
        candidate = f"{current} {word}"
        if pdfmetrics.stringWidth(candidate, font_name, font_size) <= max_width:
            current = candidate
        else:
            lines.append(current)
            current = word
    lines.append(current)
    return lines


def excel_color_to_reportlab(color):
    if not color or color.type != "rgb" or not color.rgb:
        return None


EXCEL_OPERATORS = {
    ast.Add: operator.add,
    ast.Sub: operator.sub,
    ast.Mult: operator.mul,
    ast.Div: operator.truediv,
    ast.Pow: operator.pow,
    ast.USub: operator.neg,
    ast.UAdd: operator.pos,
}


def safe_numeric_expression(expression: str):
    def evaluate(node):
        if isinstance(node, ast.Expression):
            return evaluate(node.body)
        if isinstance(node, ast.Constant) and isinstance(node.value, (int, float)):
            return node.value
        if isinstance(node, ast.BinOp) and type(node.op) in EXCEL_OPERATORS:
            return EXCEL_OPERATORS[type(node.op)](evaluate(node.left), evaluate(node.right))
        if isinstance(node, ast.UnaryOp) and type(node.op) in EXCEL_OPERATORS:
            return EXCEL_OPERATORS[type(node.op)](evaluate(node.operand))
        raise ValueError("Desteklenmeyen formül")
    return evaluate(ast.parse(expression, mode="eval"))


def resolve_excel_value(sheet, cell, cache, stack=None):
    if cell.coordinate in cache:
        return cache[cell.coordinate]
    value = cell.value
    if not isinstance(value, str) or not value.startswith("="):
        cache[cell.coordinate] = value
        return value

    stack = set(stack or ())
    if cell.coordinate in stack:
        return value
    stack.add(cell.coordinate)
    expression = value[1:].upper()

    def replace_sum(match):
        start, end = match.group(1), match.group(2)
        total = 0
        for row in sheet[start:end]:
            for item in row:
                resolved = resolve_excel_value(sheet, item, cache, stack)
                if isinstance(resolved, (int, float)):
                    total += resolved
        return str(total)

    expression = re.sub(r"SUM\(([A-Z]+\d+):([A-Z]+\d+)\)", replace_sum, expression)

    def replace_reference(match):
        referenced = resolve_excel_value(sheet, sheet[match.group(0)], cache, stack)
        return str(referenced if isinstance(referenced, (int, float)) else 0)

    expression = re.sub(r"\b[A-Z]+\d+\b", replace_reference, expression)
    try:
        result = safe_numeric_expression(expression)
    except Exception:
        result = value
    cache[cell.coordinate] = result
    return result


def format_excel_display(value, number_format: str) -> str:
    if value is None:
        return ""
    if isinstance(value, (int, float)) and "₺" in (number_format or ""):
        formatted = f"{value:,.2f}".replace(",", "_").replace(".", ",").replace("_", ".")
        return f"TL {formatted}"
    if isinstance(value, float):
        return f"{value:g}"
    return str(value)
    value = color.rgb[-6:]
    try:
        return colors.HexColor(f"#{value}")
    except ValueError:
        return None


def create_pdf_from_xlsx_layout(input_path: Path, output_path: Path) -> None:
    workbook = load_workbook(filename=str(input_path), data_only=False)
    font_name = register_unicode_pdf_font()
    visible_sheets = [sheet for sheet in workbook.worksheets if sheet.sheet_state == "visible"]
    first_orientation = visible_sheets[0].page_setup.orientation if visible_sheets else None
    page_size = landscape(A4) if first_orientation == "landscape" else A4
    document = SimpleDocTemplate(
        str(output_path),
        pagesize=page_size,
        leftMargin=24,
        rightMargin=24,
        topMargin=28,
        bottomMargin=28,
    )
    available_width = page_size[0] - 48
    story = []

    try:
        for sheet_index, sheet in enumerate(visible_sheets):
            max_row = sheet.max_row
            max_column = sheet.max_column
            while max_row > 1 and all(sheet.cell(max_row, column).value is None for column in range(1, max_column + 1)):
                max_row -= 1
            while max_column > 1 and all(sheet.cell(row, max_column).value is None for row in range(1, max_row + 1)):
                max_column -= 1

            raw_widths = []
            for column in range(1, max_column + 1):
                letter = get_column_letter(column)
                excel_width = sheet.column_dimensions[letter].width or 10
                raw_widths.append(max(28, min(180, excel_width * 6.5)))
            width_scale = min(1, available_width / sum(raw_widths))
            column_widths = [width * width_scale for width in raw_widths]

            body_style = ParagraphStyle(
                "ExcelCell",
                fontName=font_name,
                fontSize=7.5,
                leading=9,
                textColor=colors.black,
                wordWrap="CJK",
            )
            data = []
            formula_cache = {}
            for row in range(1, max_row + 1):
                output_row = []
                for column in range(1, max_column + 1):
                    cell = sheet.cell(row, column)
                    value = format_excel_display(
                        resolve_excel_value(sheet, cell, formula_cache),
                        cell.number_format,
                    )
                    markup = escape(value).replace("\n", "<br/>")
                    if cell.font.bold:
                        markup = f"<b>{markup}</b>"
                    output_row.append(Paragraph(markup, body_style))
                data.append(output_row)

            row_heights = [
                max(14, min(100, sheet.row_dimensions[row].height or 15))
                for row in range(1, max_row + 1)
            ]
            table = Table(
                data,
                colWidths=column_widths,
                rowHeights=row_heights,
                repeatRows=1 if max_row > 1 else 0,
                splitByRow=True,
            )
            commands = [
                ("FONTNAME", (0, 0), (-1, -1), font_name),
                ("FONTSIZE", (0, 0), (-1, -1), 7.5),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#BFC0CA")),
                ("LEFTPADDING", (0, 0), (-1, -1), 3),
                ("RIGHTPADDING", (0, 0), (-1, -1), 3),
                ("TOPPADDING", (0, 0), (-1, -1), 2),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 2),
            ]
            for row in range(1, max_row + 1):
                for column in range(1, max_column + 1):
                    cell = sheet.cell(row, column)
                    position = (column - 1, row - 1)
                    fill = excel_color_to_reportlab(cell.fill.fgColor) if cell.fill.fill_type else None
                    font_color = excel_color_to_reportlab(cell.font.color)
                    if fill:
                        commands.append(("BACKGROUND", position, position, fill))
                    if font_color:
                        commands.append(("TEXTCOLOR", position, position, font_color))
                    if cell.font.bold:
                        commands.append(("FONTNAME", position, position, font_name))
                    alignment = (cell.alignment.horizontal or "left").upper()
                    if alignment in {"LEFT", "CENTER", "RIGHT", "JUSTIFY"}:
                        commands.append(("ALIGN", position, position, alignment))

            for merged_range in sheet.merged_cells.ranges:
                if merged_range.max_row <= max_row and merged_range.max_col <= max_column:
                    commands.append((
                        "SPAN",
                        (merged_range.min_col - 1, merged_range.min_row - 1),
                        (merged_range.max_col - 1, merged_range.max_row - 1),
                    ))
            table.setStyle(TableStyle(commands))
            story.append(table)
            if sheet_index < len(visible_sheets) - 1:
                story.extend([PageBreak(), Spacer(1, 1)])
        document.build(story)
    finally:
        workbook.close()


def find_libreoffice_executable():
    candidates = [
        shutil.which("soffice"),
        shutil.which("libreoffice"),
        r"C:\Program Files\LibreOffice\program\soffice.com",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.com",
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return str(candidate)
    return None


def create_pdf_with_libreoffice(input_path: Path, output_path: Path) -> bool:
    executable = find_libreoffice_executable()
    if not executable:
        return False

    profile_dir = input_path.parent / "libreoffice-profile"
    profile_dir.mkdir(exist_ok=True)
    command = [
        executable,
        f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(input_path.parent),
        str(input_path),
    ]
    try:
        result = subprocess.run(
            command,
            capture_output=True,
            text=True,
            timeout=180,
            check=False,
        )
    except (OSError, subprocess.SubprocessError):
        logging.exception("LibreOffice conversion could not be started")
        return False

    generated_path = input_path.with_suffix(".pdf")
    if result.returncode != 0 or not generated_path.exists():
        logging.error(
            "LibreOffice conversion failed (%s): %s %s",
            result.returncode,
            result.stdout,
            result.stderr,
        )
        return False
    if output_path.exists():
        output_path.unlink()
    generated_path.replace(output_path)
    return True


def create_pdf_from_dwg(input_path: Path, output_path: Path) -> None:
    """Use a configured CAD engine or the free LibreDWG SVG pipeline."""
    template = os.environ.get("DWG_TO_PDF_COMMAND", "").strip()
    if template:
        if "{input}" not in template or "{output}" not in template:
            logging.error("DWG_TO_PDF_COMMAND must contain {input} and {output}")
            raise ConversionServiceUnavailable("DWG dönüştürme motoru yapılandırması geçersiz.")
        try:
            tokens = shlex.split(template, posix=os.name != "nt")
            command = [token.replace("{input}", str(input_path)).replace("{output}", str(output_path)) for token in tokens]
            result = subprocess.run(command, capture_output=True, text=True, timeout=180, check=False, cwd=input_path.parent)
        except (OSError, ValueError, subprocess.SubprocessError) as exc:
            logging.exception("DWG conversion could not be started")
            raise ConversionServiceUnavailable("DWG dönüştürme motoru çalıştırılamadı.") from exc
        if result.returncode != 0 or not output_path.exists():
            logging.error("DWG conversion failed (%s): %s %s", result.returncode, result.stdout, result.stderr)
            raise ValueError("DWG dosyası PDF'e dönüştürülemedi; çizim bozuk veya desteklenmeyen bir sürüm olabilir.")
    else:
        create_pdf_from_dwg_with_libredwg(input_path, output_path)
    validate_pdf_bytes(output_path.read_bytes())


def find_libredwg_svg_executable():
    configured = os.environ.get("LIBREDWG_DWG2SVG", "").strip()
    project_bin = Path(__file__).resolve().parent / "tools" / "libredwg" / (
        "dwg2SVG.exe" if os.name == "nt" else "dwg2SVG"
    )
    candidates = [configured, shutil.which("dwg2SVG"), shutil.which("dwg2svg"), project_bin]
    for candidate in candidates:
        if candidate and Path(candidate).is_file():
            return str(candidate)
    return None


def find_libredwg_dxf_executable():
    configured = os.environ.get("LIBREDWG_DWG2DXF", "").strip()
    project_bin = Path(__file__).resolve().parent / "tools" / "libredwg" / (
        "dwg2dxf.exe" if os.name == "nt" else "dwg2dxf"
    )
    candidates = [configured, shutil.which("dwg2dxf"), project_bin]
    for candidate in candidates:
        if candidate and Path(candidate).is_file():
            return str(candidate)
    return None


def create_pdf_from_dwg_with_libredwg(input_path: Path, output_path: Path) -> None:
    dxf_executable = find_libredwg_dxf_executable()
    executable = find_libredwg_svg_executable()
    if not executable:
        if dxf_executable:
            create_pdf_from_dwg_via_dxf(dxf_executable, input_path, output_path)
            return
        raise ConversionServiceUnavailable("Ücretsiz LibreDWG motoru sunucuda bulunamadı.")

    svg_path = output_path.with_suffix(".svg")
    def run_libredwg(model_space=False):
        command = [executable]
        if model_space:
            command.append("--mspace")
        command.append(str(input_path))
        return subprocess.run(
            command, capture_output=True, timeout=180, check=False, cwd=input_path.parent
        )

    try:
        result = run_libredwg()
        svg_prefix = result.stdout.split(b"<defs>", 1)[0].lower()
        visible_elements = (b"<path ", b"<use ", b"<text ", b"<circle ", b"<ellipse ", b"<line ", b"<polyline ")
        if result.returncode == 0 and not any(element in svg_prefix for element in visible_elements):
            logging.info("LibreDWG paper space is empty; retrying DWG in model space")
            result = run_libredwg(model_space=True)
    except (OSError, subprocess.SubprocessError) as exc:
        logging.exception("LibreDWG could not be started")
        raise ConversionServiceUnavailable("LibreDWG motoru çalıştırılamadı.") from exc
    if result.returncode != 0 or b"<svg" not in result.stdout[:4096].lower():
        logging.error("LibreDWG failed (%s): %s", result.returncode, result.stderr.decode("utf-8", "replace"))
        raise ValueError("LibreDWG bu DWG çizimini okuyamadı veya desteklenen geometri bulamadı.")
    if b"<use " in result.stdout.lower() and dxf_executable:
        logging.info("LibreDWG output contains block references; switching to the DXF renderer")
        create_pdf_from_dwg_via_dxf(dxf_executable, input_path, output_path)
        return
    try:
        from svglib.svglib import svg2rlg
        from reportlab.graphics import renderPDF

        svg_data = result.stdout
        viewbox_match = re.search(rb'viewBox="([^"]+)"', svg_data)
        if viewbox_match:
            viewbox = [float(value) for value in viewbox_match.group(1).split()]
            if len(viewbox) == 4:
                x, y, width, height = viewbox
                if abs(x) > max(width, 1) * 0.5 or abs(y) > max(height, 1) * 0.5:
                    seed = max(abs(value) for value in viewbox) * 2 + max(width, height) + 1
                    probe_viewbox = f'viewBox="0 0 {seed:.6f} {seed:.6f}"'.encode()
                    probe_data = svg_data[:viewbox_match.start()] + probe_viewbox + svg_data[viewbox_match.end():]
                    svg_path.write_bytes(probe_data)
                    probe = svg2rlg(str(svg_path))
                    if probe is not None:
                        left, bottom, right, top = probe.getBounds()
                        actual_top = seed - top
                        actual_bottom = seed - bottom
                        drawing_width = max(right - left, 1.0)
                        drawing_height = max(actual_bottom - actual_top, 1.0)
                        padding = max(drawing_width, drawing_height) * 0.01
                        corrected = (
                            f'viewBox="{left - padding:.6f} {actual_top - padding:.6f} '
                            f'{drawing_width + padding * 2:.6f} {drawing_height + padding * 2:.6f}"'
                        ).encode()
                        corrected_match = re.search(rb'viewBox="[^"]+"', probe_data)
                        svg_data = (
                            probe_data[:corrected_match.start()]
                            + corrected
                            + probe_data[corrected_match.end():]
                        )
                        logging.info("Corrected invalid LibreDWG SVG viewBox from drawing bounds")
        svg_path.write_bytes(svg_data)
        drawing = svg2rlg(str(svg_path))
        if drawing is None:
            raise ValueError("SVG çizimi okunamadı.")
        renderPDF.drawToFile(drawing, str(output_path))
    except (ImportError, OSError) as exc:
        logging.exception("SVG to PDF renderer is unavailable")
        raise ConversionServiceUnavailable("SVG → PDF bileşeni sunucuda kullanılamıyor.") from exc
    except Exception as exc:
        logging.exception("LibreDWG SVG output could not be rendered")
        raise ValueError("DWG çizimi PDF sayfasına aktarılamadı.") from exc


def create_pdf_from_dwg_via_dxf(executable: str, input_path: Path, output_path: Path) -> None:
    dxf_path = output_path.with_suffix(".dxf")
    try:
        result = subprocess.run(
            [executable, "-y", "-o", str(dxf_path), str(input_path)],
            capture_output=True,
            timeout=180,
            check=False,
            cwd=input_path.parent,
        )
    except (OSError, subprocess.SubprocessError) as exc:
        logging.exception("LibreDWG DXF conversion could not be started")
        raise ConversionServiceUnavailable("LibreDWG DXF motoru çalıştırılamadı.") from exc
    if result.returncode != 0 or not dxf_path.exists():
        logging.error("LibreDWG DXF conversion failed (%s): %s", result.returncode, result.stderr.decode("utf-8", "replace"))
        raise ValueError("LibreDWG bu DWG çizimini DXF'e aktaramadı.")

    try:
        cache_root = output_path.parent / "cad-render-cache"
        cache_root.mkdir(exist_ok=True)
        os.environ.setdefault("MPLCONFIGDIR", str(cache_root / "matplotlib"))
        os.environ.setdefault("XDG_CACHE_HOME", str(cache_root))

        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import ezdxf
        from ezdxf.addons.drawing import RenderContext, Frontend
        from ezdxf.addons.drawing.config import Configuration, BackgroundPolicy
        from ezdxf.addons.drawing.matplotlib import MatplotlibBackend

        document = ezdxf.readfile(dxf_path)
        modelspace = document.modelspace()
        if not len(modelspace):
            raise ValueError("DXF model alanı boş.")
        figure = plt.figure(frameon=False)
        try:
            axes = figure.add_axes([0, 0, 1, 1])
            axes.set_axis_off()
            context = RenderContext(document)
            context.set_current_layout(modelspace)
            backend = MatplotlibBackend(axes)
            configuration = Configuration(background_policy=BackgroundPolicy.WHITE)
            Frontend(context, backend, config=configuration).draw_layout(modelspace, finalize=True)
            figure.savefig(
                output_path,
                format="pdf",
                bbox_inches="tight",
                pad_inches=0.05,
                facecolor="white",
            )
        finally:
            plt.close(figure)
    except (ImportError, OSError) as exc:
        logging.exception("DXF to PDF renderer is unavailable")
        raise ConversionServiceUnavailable("DXF → PDF bileşeni sunucuda kullanılamıyor.") from exc
    except Exception as exc:
        logging.exception("DXF output could not be rendered")
        raise ValueError("DWG çizimi PDF sayfasına aktarılamadı.") from exc
    validate_pdf_bytes(output_path.read_bytes())


def attachment_response(data: bytes, filename: str, mimetype: str):
    return send_file(
        io.BytesIO(data),
        as_attachment=True,
        download_name=filename,
        mimetype=mimetype,
    )


@app.get("/")
def index():
    index_path = Path(__file__).resolve().parent / "index.html"
    if index_path.exists():
        return send_file(index_path, mimetype="text/html")

    return jsonify({
        "service": "pdf-converter-api",
        "supported_operations": sorted(SUPPORTED_OPERATIONS.keys()),
    })


@app.get("/health")
def health():
    return jsonify({"ok": True, "service": "pdf-converter-api", "version": APP_VERSION})


@app.post("/api/convert")
@limiter.limit("10 per minute")
def convert():
    file = request.files.get("file")
    operation = request.form.get("operation", "")

    if not file or not file.filename:
        return jsonify({"error": "Dosya yüklenmedi."}), 400
    if not operation:
        return jsonify({"error": "Dönüştürme türü belirtilmedi."}), 400

    extension = Path(file.filename).suffix.lower()
    output_suffix = {
        "pdf-to-docx": ".docx",
        "pdf-to-xlsx": ".xlsx",
        "pdf-to-pptx": ".pptx",
        "docx-to-pdf": ".pdf",
        "xlsx-to-pdf": ".pdf",
        "pptx-to-pdf": ".pdf",
        "dwg-to-pdf": ".pdf",
    }.get(operation)

    if not output_suffix:
        return jsonify({"error": "Desteklenmeyen işlem."}), 400

    allowed_extensions = SUPPORTED_OPERATIONS.get(operation, set())
    if extension not in allowed_extensions:
        expected = ", ".join(sorted(allowed_extensions)) or "yok"
        return jsonify({"error": f"Bu işlem için {extension or 'dosya uzantısı'} desteklenmiyor. Beklenen: {expected}"}), 400

    with tempfile.TemporaryDirectory() as tmpdir:
        temp_dir = Path(tmpdir)
        input_path = temp_dir / f"input{extension}"
        output_path = temp_dir / f"converted{output_suffix}"
        file.save(input_path)

        try:
            if extension == ".pdf":
                validate_pdf_bytes(input_path.read_bytes())
                validation_reader = PdfReader(str(input_path), strict=False)
                if validation_reader.is_encrypted:
                    return jsonify({"error": "Dönüştürmeden önce PDF kilidini açın."}), 400
                enforce_pdf_page_limit(validation_reader)
            elif extension in OFFICE_MARKERS:
                validate_office_archive(input_path, extension)
            elif extension == ".dwg":
                validate_dwg_file(input_path)
            if operation == "pdf-to-docx":
                text = extract_text_from_pdf(input_path)
                word_mode = request.form.get("pdf_word_mode", "editable")
                if word_mode == "layout":
                    create_docx_from_scanned_pdf(input_path, output_path)
                elif word_mode == "hybrid":
                    create_hybrid_docx(input_path, output_path)
                elif word_mode != "editable":
                    return jsonify({"error": "Geçersiz Word aktarım biçimi."}), 400
                elif text.strip():
                    create_docx_from_text(text, output_path)
                else:
                    if not create_editable_docx_with_ocr(input_path, output_path):
                        create_docx_from_scanned_pdf(input_path, output_path)
            elif operation == "pdf-to-xlsx":
                if not create_xlsx_from_pdf_tables(input_path, output_path):
                    text = extract_text_from_pdf(input_path)
                    create_xlsx_from_text(text, output_path)
            elif operation == "pdf-to-pptx":
                text = extract_text_from_pdf(input_path)
                create_pptx_from_text(text, output_path)
            elif operation == "docx-to-pdf":
                text = extract_text_from_docx(input_path)
                create_pdf_from_text(text, output_path)
            elif operation == "xlsx-to-pdf":
                if not create_pdf_with_libreoffice(input_path, output_path):
                    create_pdf_from_xlsx_layout(input_path, output_path)
            elif operation == "pptx-to-pdf":
                text = extract_text_from_pptx(input_path)
                create_pdf_from_text(text, output_path)
            elif operation == "dwg-to-pdf":
                create_pdf_from_dwg(input_path, output_path)
            else:
                return jsonify({"error": "Desteklenmeyen işlem."}), 400
        except ConversionServiceUnavailable as exc:
            return jsonify({"error": str(exc)}), 503
        except ValueError as exc:
            return jsonify({"error": str(exc)}), 400
        except Exception:  # pragma: no cover - runtime safety
            logging.exception("Conversion failed")
            return jsonify({"error": "Dönüştürme sırasında güvenli biçimde işlenemeyen bir dosyayla karşılaşıldı."}), 500

        output_data = output_path.read_bytes()

    original_stem = Path(file.filename).stem.strip()
    response_stem = re.sub(r'[<>:"/\\|?*\x00-\x1F]', "_", original_stem) or "donusturulmus"
    response_name = f"{response_stem}{output_suffix}"
    mime_types = {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    return send_file(
        io.BytesIO(output_data),
        as_attachment=True,
        download_name=response_name,
        mimetype=mime_types[output_suffix],
    )


@app.post("/api/pdf-tool")
@limiter.limit("15 per minute")
def pdf_tool():
    uploaded = request.files.get("file")
    operation = request.form.get("operation", "")
    if not uploaded or not uploaded.filename:
        return jsonify({"error": "PDF dosyası seçilmedi."}), 400
    if Path(uploaded.filename).suffix.lower() != ".pdf":
        return jsonify({"error": "Bu araç yalnızca PDF dosyalarını kabul eder."}), 400

    supported = {"protect", "unlock", "repair", "compare"}
    if operation not in supported:
        return jsonify({"error": "Desteklenmeyen PDF aracı."}), 400

    try:
        source_bytes = uploaded.read()
        validate_pdf_bytes(source_bytes)
        reader = PdfReader(io.BytesIO(source_bytes), strict=False)
        password = request.form.get("password", "")

        if operation == "unlock":
            if not reader.is_encrypted:
                enforce_pdf_page_limit(reader)
                return attachment_response(source_bytes, "kilitsiz.pdf", "application/pdf")
            if not password or reader.decrypt(password) == 0:
                return jsonify({"error": "PDF parolası yanlış."}), 400
            enforce_pdf_page_limit(reader)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            output = io.BytesIO()
            writer.write(output)
            return attachment_response(output.getvalue(), "kilitsiz.pdf", "application/pdf")

        if reader.is_encrypted:
            return jsonify({"error": "Bu işlemden önce PDF kilidini açın."}), 400
        enforce_pdf_page_limit(reader)

        if operation == "protect":
            if len(password) < 4:
                return jsonify({"error": "Parola en az 4 karakter olmalıdır."}), 400
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            writer.encrypt(user_password=password, owner_password=password)
            output = io.BytesIO()
            writer.write(output)
            return attachment_response(output.getvalue(), "parolali.pdf", "application/pdf")

        if operation == "repair":
            writer = PdfWriter()
            recovered = 0
            for page in reader.pages:
                try:
                    writer.add_page(page)
                    recovered += 1
                except Exception:
                    logging.exception("A damaged PDF page could not be recovered")
            if not recovered:
                return jsonify({"error": "Kurtarılabilecek bir sayfa bulunamadı."}), 422
            output = io.BytesIO()
            writer.write(output)
            return attachment_response(output.getvalue(), "onarilmis.pdf", "application/pdf")

        second = request.files.get("second_file")
        if not second or not second.filename or Path(second.filename).suffix.lower() != ".pdf":
            return jsonify({"error": "Karşılaştırma için ikinci PDF gereklidir."}), 400
        second_bytes = second.read()
        validate_pdf_bytes(second_bytes)
        second_reader = PdfReader(io.BytesIO(second_bytes), strict=False)
        if second_reader.is_encrypted:
            return jsonify({"error": "İkinci PDF parola korumalı olmamalıdır."}), 400
        enforce_pdf_page_limit(second_reader)
        first_lines = extract_text_from_reader(reader).splitlines()
        second_lines = extract_text_from_reader(second_reader).splitlines()
        diff = "\n".join(difflib.unified_diff(
            first_lines,
            second_lines,
            fromfile=secure_filename(uploaded.filename),
            tofile=secure_filename(second.filename),
            lineterm="",
        ))
        report = diff or "Metin içeriğinde fark bulunamadı."
        return attachment_response(report.encode("utf-8"), "karsilastirma.txt", "text/plain; charset=utf-8")
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    except Exception:
        logging.exception("PDF tool failed")
        return jsonify({"error": "PDF güvenli biçimde işlenemedi."}), 500


@app.post("/api/create-document")
@limiter.limit("20 per minute")
def create_document():
    payload = request.get_json(silent=True) or {}
    document_type = payload.get("type")
    content = payload.get("content")
    requested_name = str(payload.get("filename") or "yeni-belge")
    extensions = {"word": ".docx", "excel": ".xlsx", "powerpoint": ".pptx"}
    extension = extensions.get(document_type)
    if not extension or not isinstance(content, list):
        return jsonify({"error": "Geçersiz belge oluşturma isteği."}), 400
    base_name = Path(requested_name).stem.strip()
    safe_base = re.sub(r'[<>:"/\\|?*\x00-\x1F]', "-", base_name)[:80] or "yeni-belge"
    output_name = f"{safe_base}{extension}"
    output = io.BytesIO()

    try:
        if document_type == "word":
            document = Document()
            for item in content[:1000]:
                text = clean_document_text(str(item.get("text", "")))[:20000]
                if not text:
                    continue
                if item.get("type") == "heading":
                    document.add_heading(text, level=1)
                else:
                    document.add_paragraph(text)
            document.save(output)
            mimetype = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        elif document_type == "excel":
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "Sayfa1"
            for row in content[:2000]:
                values = []
                for value in list(row)[:100]:
                    cleaned = clean_document_text(str(value))[:32767]
                    if cleaned.startswith("="):
                        values.append(cleaned)
                    else:
                        values.append(parse_turkish_number(cleaned))
                sheet.append(values)
            sheet.freeze_panes = "A2"
            for column in range(1, min(sheet.max_column, 100) + 1):
                values = [str(sheet.cell(row, column).value or "") for row in range(1, min(sheet.max_row, 100) + 1)]
                sheet.column_dimensions[get_column_letter(column)].width = min(40, max(10, max(map(len, values), default=8) + 2))
            workbook.save(output)
            mimetype = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        else:
            presentation = Presentation()
            for slide_data in content[:200]:
                slide = presentation.slides.add_slide(presentation.slide_layouts[1])
                slide.shapes.title.text = clean_document_text(str(slide_data.get("title", "")))[:500]
                body = clean_document_text(str(slide_data.get("body", "")))[:10000]
                frame = slide.placeholders[1].text_frame
                frame.clear()
                for index, line in enumerate(body.splitlines() or [""]):
                    paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                    paragraph.text = line
            presentation.save(output)
            mimetype = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    except Exception as exc:
        logging.exception("Document creation failed")
        return jsonify({"error": f"Belge oluşturulamadı: {exc}"}), 500

    output.seek(0)
    return send_file(
        output,
        as_attachment=True,
        download_name=output_name,
        mimetype=mimetype,
    )


def extract_text_from_reader(reader: PdfReader) -> str:
    return "\n".join((page.extract_text() or "") for page in reader.pages)


if __name__ == "__main__":
    app.run(
        host=os.environ.get("HOST", "0.0.0.0"),
        port=int(os.environ.get("PORT", "3000")),
        debug=os.environ.get("FLASK_DEBUG") == "1",
    )
